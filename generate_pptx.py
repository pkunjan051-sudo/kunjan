import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Define Color Palette (Netflix/IMDb Cinematic Dark Theme)
BG_COLOR = RGBColor(13, 14, 21)        # #0D0E15 Obsidian Black
CARD_BG = RGBColor(23, 24, 34)         # #171822 Dark Charcoal Surface
CARD_BORDER = RGBColor(42, 44, 61)     # #2A2C3D Subtle Slate Border
CRIMSON = RGBColor(229, 9, 20)         # #E50914 Netflix Crimson Red
RED_GLOW = RGBColor(255, 46, 76)       # #FF2E4C Accent Glow
GOLD = RGBColor(255, 193, 7)           # #FFC107 Rating Gold
TEXT_WHITE = RGBColor(255, 255, 255)   # Primary Text
TEXT_GRAY = RGBColor(185, 188, 208)   # High-Contrast Secondary Text for Projectors
TEXT_MUTED = RGBColor(130, 134, 160)  # Muted Captions
BLUE_ACCENT = RGBColor(56, 189, 248)  # Tech Blue

def create_presentation():
    prs = Presentation()
    # Set 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    def set_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_header(slide, title_text, category_text="CINEMACENTRAL • FLUTTER ENTERTAINMENT APP"):
        # Header Badge / Category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.733), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(12)
        p_cat.font.bold = True
        p_cat.font.color.rgb = CRIMSON
        p_cat.font.name = "Arial"

        # Main Slide Title (34pt for Projector Distance Visibility)
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.733), Inches(0.85))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(34)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE
        p_title.font.name = "Arial"

    def add_card(slide, left, top, width, height, bg_rgb=CARD_BG, border_rgb=CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_rgb
        if border_rgb:
            shape.line.color.rgb = border_rgb
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    def add_flow_diagram(slide, left, top, width, height, steps):
        n = len(steps)
        step_width = (width - Inches(0.35 * (n - 1))) / n
        for i, step in enumerate(steps):
            s_left = left + i * (step_width + Inches(0.35))
            # Card
            add_card(slide, s_left, top, step_width, height, CARD_BG, CRIMSON if i == n-1 else CARD_BORDER)
            # Text inside step
            tb = slide.shapes.add_textbox(s_left + Inches(0.12), top + Inches(0.12), step_width - Inches(0.24), height - Inches(0.24))
            tf = tb.text_frame
            tf.word_wrap = True
            
            p0 = tf.paragraphs[0]
            p0.text = f"STEP 0{i+1}"
            p0.font.size = Pt(12)
            p0.font.bold = True
            p0.font.color.rgb = CRIMSON if i == n-1 else GOLD
            p0.alignment = PP_ALIGN.CENTER
            
            p1 = tf.add_paragraph()
            p1.text = step["title"]
            p1.font.size = Pt(16)
            p1.font.bold = True
            p1.font.color.rgb = TEXT_WHITE
            p1.alignment = PP_ALIGN.CENTER
            
            if "desc" in step:
                p2 = tf.add_paragraph()
                p2.text = step["desc"]
                p2.font.size = Pt(13.5)
                p2.font.color.rgb = TEXT_GRAY
                p2.alignment = PP_ALIGN.CENTER

            # Arrow between steps
            if i < n - 1:
                arr_box = slide.shapes.add_textbox(s_left + step_width, top + height/2 - Inches(0.35), Inches(0.35), Inches(0.7))
                tf_arr = arr_box.text_frame
                p_arr = tf_arr.paragraphs[0]
                p_arr.text = "➔"
                p_arr.font.size = Pt(24)
                p_arr.font.bold = True
                p_arr.font.color.rgb = CRIMSON
                p_arr.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide1)

    # Hero Accent Card Background
    add_card(slide1, Inches(0.8), Inches(0.7), Inches(11.733), Inches(6.0), CARD_BG, CRIMSON)

    # App Badge
    badge_shape = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(1.2), Inches(3.2), Inches(0.55))
    badge_shape.fill.solid()
    badge_shape.fill.fore_color.rgb = CRIMSON
    badge_shape.line.fill.background()
    tf_b = badge_shape.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "🎬 FLUTTER MOBILE APP"
    p_b.font.size = Pt(13)
    p_b.font.bold = True
    p_b.font.color.rgb = TEXT_WHITE
    p_b.alignment = PP_ALIGN.CENTER

    # Title Text
    tb1 = slide1.shapes.add_textbox(Inches(1.3), Inches(2.0), Inches(10.7), Inches(2.7))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "CINEMACENTRAL"
    p1.font.size = Pt(50)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE

    p2 = tf1.add_paragraph()
    p2.text = "Next-Generation Movie Discovery & Entertainment Platform"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = RED_GLOW

    p3 = tf1.add_paragraph()
    p3.text = "A high-performance mobile application engineered with Flutter, Dart, REST APIs, and local persistence."
    p3.font.size = Pt(16)
    p3.font.color.rgb = TEXT_GRAY

    # Footer Card Info
    add_card(slide1, Inches(1.3), Inches(4.9), Inches(10.733), Inches(1.5), BG_COLOR, CARD_BORDER)
    tb_foot = slide1.shapes.add_textbox(Inches(1.5), Inches(5.05), Inches(10.333), Inches(1.2))
    tf_f = tb_foot.text_frame
    tf_f.word_wrap = True
    
    pf1 = tf_f.paragraphs[0]
    pf1.text = "PROJECT METADATA & TECHNICAL STACK"
    pf1.font.size = Pt(13)
    pf1.font.bold = True
    pf1.font.color.rgb = GOLD

    pf2 = tf_f.add_paragraph()
    pf2.text = "• Technology: Flutter 3.x  |  Dart SDK  |  TMDB REST API v3  |  HTTP Client  |  SharedPreferences\n• Topic: Mobile Software Engineering & UI/UX Design  |  College Project Viva Presentation"
    pf2.font.size = Pt(15)
    pf2.font.color.rgb = TEXT_WHITE

    # ==========================================
    # SLIDE 2: Introduction
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide2)
    add_header(slide2, "Introduction: What is Cinemacentral?")

    # Left Column: Overview Card
    add_card(slide2, Inches(0.8), Inches(1.65), Inches(5.7), Inches(5.15))
    tb_intro1 = slide2.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(5.3), Inches(4.75))
    tf_i1 = tb_intro1.text_frame
    tf_i1.word_wrap = True

    pi1 = tf_i1.paragraphs[0]
    pi1.text = "Application Concept & Vision"
    pi1.font.size = Pt(22)
    pi1.font.bold = True
    pi1.font.color.rgb = CRIMSON

    pi2 = tf_i1.add_paragraph()
    pi2.text = "\nCinemacentral is a modern cross-platform mobile discovery application built with Flutter and Dart.\n\nIt serves as a single unified hub for exploring real-time trending movies, top-rated films, upcoming releases, short video trailers, actor filmographies, and persistent watchlists."
    pi2.font.size = Pt(15)
    pi2.font.color.rgb = TEXT_GRAY

    # Right Column: 3 Purpose Cards
    purposes = [
        {"title": "🎯 Centralized Discovery", "desc": "Consolidates global cinema data into a single modern interface without switching platforms."},
        {"title": "⚡ Immersive UX & Reels", "desc": "Features vertical video trailer reels and poster carousels engineered for fluid touch interaction."},
        {"title": "📌 Local Watchlist", "desc": "Enables instant local persistence via SharedPreferences to save movies for later viewing."}
    ]
    for idx, purp in enumerate(purposes):
        top_pos = Inches(1.65 + idx * 1.75)
        add_card(slide2, Inches(6.8), top_pos, Inches(5.733), Inches(1.6))
        tb_p = slide2.shapes.add_textbox(Inches(7.0), top_pos + Inches(0.15), Inches(5.333), Inches(1.3))
        tf_p = tb_p.text_frame
        tf_p.word_wrap = True
        
        pp1 = tf_p.paragraphs[0]
        pp1.text = purp["title"]
        pp1.font.size = Pt(18)
        pp1.font.bold = True
        pp1.font.color.rgb = TEXT_WHITE
        
        pp2 = tf_p.add_paragraph()
        pp2.text = purp["desc"]
        pp2.font.size = Pt(14)
        pp2.font.color.rgb = TEXT_GRAY

    # ==========================================
    # SLIDE 3: Problem Statement
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide3)
    add_header(slide3, "Problem Statement: Challenges in Existing Solutions")

    problems = [
        {
            "num": "01",
            "title": "Fragmented Data",
            "desc": "Users must navigate between separate websites to find storylines, cast lists, ratings, and video trailers."
        },
        {
            "num": "02",
            "title": "Outdated UI/UX",
            "desc": "Traditional movie tools feature cluttered desktop layouts, uninspired fonts, and poor mobile responsiveness."
        },
        {
            "num": "03",
            "title": "Sluggish Mobile Performance",
            "desc": "Heavy web views and unoptimized image loading cause lag, buffer delays, and high battery consumption."
        }
    ]

    for i, prob in enumerate(problems):
        left_pos = Inches(0.8 + i * 3.98)
        add_card(slide3, left_pos, Inches(1.7), Inches(3.78), Inches(5.1), CARD_BG, CRIMSON if i == 0 else CARD_BORDER)
        
        tb_pr = slide3.shapes.add_textbox(left_pos + Inches(0.2), Inches(1.9), Inches(3.38), Inches(4.7))
        tf_pr = tb_pr.text_frame
        tf_pr.word_wrap = True

        p_num = tf_pr.paragraphs[0]
        p_num.text = f"CHALLENGE {prob['num']}"
        p_num.font.size = Pt(13)
        p_num.font.bold = True
        p_num.font.color.rgb = CRIMSON if i == 0 else GOLD

        p_t = tf_pr.add_paragraph()
        p_t.text = prob["title"]
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE

        p_d = tf_pr.add_paragraph()
        p_d.text = f"\n{prob['desc']}"
        p_d.font.size = Pt(15)
        p_d.font.color.rgb = TEXT_GRAY

    # ==========================================
    # SLIDE 4: Project Objectives
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide4)
    add_header(slide4, "Project Objectives & Key Goals")

    objectives = [
        {"title": "🎨 Modern Netflix-Inspired Theme", "desc": "Design a sleek, obsidian dark cinematic interface with custom typography and poster-first grids."},
        {"title": "🔌 Live REST API Integration", "desc": "Connect to TMDB REST API v3 using HTTP client to dynamically retrieve global movie datasets."},
        {"title": "📲 Cross-Platform Responsiveness", "desc": "Build responsive mobile layouts compatible across diverse Android and iOS display form factors."},
        {"title": "🚀 Smooth Navigation & State", "desc": "Implement Material 3 navigation bars, tab switching, and reactive watchlist state updates."}
    ]

    for idx, obj in enumerate(objectives):
        col = idx % 2
        row = idx // 2
        left_pos = Inches(0.8 + col * 5.95)
        top_pos = Inches(1.7 + row * 2.65)
        
        add_card(slide4, left_pos, top_pos, Inches(5.75), Inches(2.4))
        tb_o = slide4.shapes.add_textbox(left_pos + Inches(0.25), top_pos + Inches(0.2), Inches(5.25), Inches(2.0))
        tf_o = tb_o.text_frame
        tf_o.word_wrap = True

        po1 = tf_o.paragraphs[0]
        po1.text = obj["title"]
        po1.font.size = Pt(19)
        po1.font.bold = True
        po1.font.color.rgb = GOLD

        po2 = tf_o.add_paragraph()
        po2.text = f"\n{obj['desc']}"
        po2.font.size = Pt(15)
        po2.font.color.rgb = TEXT_GRAY

    # ==========================================
    # SLIDE 5: Technology Stack (Priority Slide)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide5)
    add_header(slide5, "Technology Stack & Libraries")

    tech_stack = [
        {"name": "Flutter 3.x", "tag": "CORE FRAMEWORK", "desc": "Google's UI toolkit for compiling fast, natively compiled cross-platform mobile apps.", "color": BLUE_ACCENT},
        {"name": "Dart SDK", "tag": "LANGUAGE", "desc": "Strongly typed language optimized for reactive client UI and asynchronous logic execution.", "color": BLUE_ACCENT},
        {"name": "TMDB REST API", "tag": "DATA BACKEND", "desc": "The Movie Database API v3 for fetching trending movies, search, genres, and cast bios.", "color": CRIMSON},
        {"name": "HTTP Package", "tag": "NETWORKING", "desc": "Handles async REST API requests, URI parameters, headers, and JSON response parsing.", "color": CRIMSON},
        {"name": "SharedPreferences", "tag": "LOCAL STORAGE", "desc": "Key-value persistent local storage used to save watchlists offline across app restarts.", "color": GOLD},
        {"name": "Google Fonts & Media", "tag": "TYPOGRAPHY & VIDEO", "desc": "Montserrat & Inter typography + VideoPlayer for vertical trailer reels.", "color": GOLD}
    ]

    for idx, tech in enumerate(tech_stack):
        col = idx % 3
        row = idx // 3
        left_pos = Inches(0.8 + col * 3.98)
        top_pos = Inches(1.7 + row * 2.65)

        add_card(slide5, left_pos, top_pos, Inches(3.78), Inches(2.45))
        tb_t = slide5.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.15), Inches(3.38), Inches(2.15))
        tf_t = tb_t.text_frame
        tf_t.word_wrap = True

        pt0 = tf_t.paragraphs[0]
        pt0.text = tech["tag"]
        pt0.font.size = Pt(11)
        pt0.font.bold = True
        pt0.font.color.rgb = tech["color"]

        pt1 = tf_t.add_paragraph()
        pt1.text = tech["name"]
        pt1.font.size = Pt(20)
        pt1.font.bold = True
        pt1.font.color.rgb = TEXT_WHITE

        pt2 = tf_t.add_paragraph()
        pt2.text = f"\n{tech['desc']}"
        pt2.font.size = Pt(14)
        pt2.font.color.rgb = TEXT_GRAY

    # ==========================================
    # SLIDE 6: System & App Architecture (Priority Slide)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide6)
    add_header(slide6, "System & Application Architecture")

    # Flowchart Diagram (Enlarged Height & Text)
    arch_steps = [
        {"title": "User Touch", "desc": "Tap card, search, or tab switch"},
        {"title": "Flutter UI Views", "desc": "HomeScreen, Search, Details"},
        {"title": "ApiService Layer", "desc": "HTTP REST client & API router"},
        {"title": "TMDB / Storage", "desc": "Remote API & SharedPreferences"},
        {"title": "State Update", "desc": "ValueNotifier UI rebuild"}
    ]
    add_flow_diagram(slide6, Inches(0.8), Inches(1.65), Inches(11.733), Inches(2.35), arch_steps)

    # Architectural Highlights Box
    add_card(slide6, Inches(0.8), Inches(4.3), Inches(11.733), Inches(2.5))
    tb_arch = slide6.shapes.add_textbox(Inches(1.0), Inches(4.45), Inches(11.333), Inches(2.2))
    tf_a = tb_arch.text_frame
    tf_a.word_wrap = True

    pa1 = tf_a.paragraphs[0]
    pa1.text = "3-Tier Clean Architecture Principles"
    pa1.font.size = Pt(20)
    pa1.font.bold = True
    pa1.font.color.rgb = GOLD

    pa2 = tf_a.add_paragraph()
    pa2.text = "• Presentation Layer: Modular Material 3 views (HomeScreen, SearchScreen, ReelsScreen, MovieDetailsScreen).\n• Service Layer: ApiService decouples network requests from UI; FavoritesService manages persistent watchlist state.\n• Data Layer: Type-safe Dart models (Movie, Actor) deserialize raw JSON with null-safety fallbacks."
    pa2.font.size = Pt(14.5)
    pa2.font.color.rgb = TEXT_WHITE

    # ==========================================
    # SLIDE 7: Main Features (Priority Slide)
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide7)
    add_header(slide7, "Main Features & Application Capabilities")

    features = [
        {"title": "🚀 Animated Splash Screen", "desc": "Sleek branding startup screen with service pre-initialization."},
        {"title": "🏠 Home Discovery Hub", "desc": "Hero backdrop banner, trending rows & genre filter chips."},
        {"title": "🎬 Short Trailers Feed (Reels)", "desc": "Full-screen vertical swipe video player with inline trailer controls."},
        {"title": "🔍 Smart Search & Filter", "desc": "Query debouncing to reduce API load + real-time genre filtering."},
        {"title": "📖 Detailed Movie Pages", "desc": "Storyline summary, ratings, release info, cast carousels, & trailer launch."},
        {"title": "⭐ Persistent Watchlist", "desc": "Bookmark favorite movies saved locally using SharedPreferences."}
    ]

    for idx, feat in enumerate(features):
        col = idx % 2
        row = idx // 2
        left_pos = Inches(0.8 + col * 5.95)
        top_pos = Inches(1.65 + row * 1.75)

        add_card(slide7, left_pos, top_pos, Inches(5.75), Inches(1.6))
        tb_f = slide7.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.15), Inches(5.35), Inches(1.3))
        tf_f2 = tb_f.text_frame
        tf_f2.word_wrap = True

        pf1 = tf_f2.paragraphs[0]
        pf1.text = feat["title"]
        pf1.font.size = Pt(18)
        pf1.font.bold = True
        pf1.font.color.rgb = TEXT_WHITE

        pf2 = tf_f2.add_paragraph()
        pf2.text = feat["desc"]
        pf2.font.size = Pt(14)
        pf2.font.color.rgb = TEXT_GRAY

    # ==========================================
    # SLIDE 8: UI/UX Design System
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide8)
    add_header(slide8, "UI/UX Design System & Aesthetics")

    # Left Column: Design Tokens
    add_card(slide8, Inches(0.8), Inches(1.65), Inches(5.7), Inches(5.15))
    tb_ui = slide8.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(5.3), Inches(4.75))
    tf_u = tb_ui.text_frame
    tf_u.word_wrap = True

    pu1 = tf_u.paragraphs[0]
    pu1.text = "Cinematic Visual Aesthetics"
    pu1.font.size = Pt(22)
    pu1.font.bold = True
    pu1.font.color.rgb = CRIMSON

    pu2 = tf_u.add_paragraph()
    pu2.text = "• Dark Obsidian Canvas (#0D0E15): Ambient background designed for poster contrast and low eye strain.\n• Crimson Accent (#E50914): Highlights call-to-action buttons, active navigation, and hero badges.\n• Gold Rating Badges (#FFC107): Visual rating indicators (e.g. ⭐ 8.5/10) for instant decision making.\n• Modern Typography: Montserrat bold display headers paired with clean Inter body copy."
    pu2.font.size = Pt(14.5)
    pu2.font.color.rgb = TEXT_WHITE

    # Right Column: Visual Layout Cards
    ui_cards = [
        {"title": "🖼️ Poster-Based Browsing", "desc": "High-aspect-ratio poster cards with hover scale animations and rating badges."},
        {"title": "✨ Glassmorphism & Elevation", "desc": "Semi-transparent dark surfaces with rounded borders and clean card margins."},
        {"title": "📱 Responsive Layouts", "desc": "Auto-adjusting grid columns and scalable carousels optimized for mobile displays."}
    ]
    for idx, uic in enumerate(ui_cards):
        top_pos = Inches(1.65 + idx * 1.75)
        add_card(slide8, Inches(6.8), top_pos, Inches(5.733), Inches(1.6))
        tb_uc = slide8.shapes.add_textbox(Inches(7.0), top_pos + Inches(0.15), Inches(5.333), Inches(1.3))
        tf_uc = tb_uc.text_frame
        tf_uc.word_wrap = True

        puc1 = tf_uc.paragraphs[0]
        puc1.text = uic["title"]
        puc1.font.size = Pt(18)
        puc1.font.bold = True
        puc1.font.color.rgb = GOLD

        puc2 = tf_uc.add_paragraph()
        puc2.text = uic["desc"]
        puc2.font.size = Pt(14)
        puc2.font.color.rgb = TEXT_GRAY

    # ==========================================
    # SLIDE 9: API Integration (Priority Slide)
    # ==========================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide9)
    add_header(slide9, "API Integration & Data Flow")

    # Diagram Pipeline
    api_steps = [
        {"title": "Flutter App", "desc": "User initiates view or search"},
        {"title": "HTTP Request", "desc": "ApiService sends async GET"},
        {"title": "TMDB API Server", "desc": "Processes query & returns JSON"},
        {"title": "JSON Model", "desc": "Movie.fromJson() parsing"},
        {"title": "Flutter UI", "desc": "Renders poster grid & details"}
    ]
    add_flow_diagram(slide9, Inches(0.8), Inches(1.65), Inches(11.733), Inches(2.35), api_steps)

    # API Technical Details Card
    add_card(slide9, Inches(0.8), Inches(4.3), Inches(11.733), Inches(2.5))
    tb_api = slide9.shapes.add_textbox(Inches(1.0), Inches(4.45), Inches(11.333), Inches(2.2))
    tf_api = tb_api.text_frame
    tf_api.word_wrap = True

    papi1 = tf_api.paragraphs[0]
    papi1.text = "Technical API Implementation & Resilience"
    papi1.font.size = Pt(20)
    papi1.font.bold = True
    papi1.font.color.rgb = CRIMSON

    papi2 = tf_api.add_paragraph()
    papi2.text = "• REST Endpoints: Serves /trending/movie/week, /movie/popular, /movie/top_rated, and /search/movie.\n• CDN Image Hosting: Constructs dynamic image URLs (https://image.tmdb.org/t/p/w500) for fast poster loading.\n• Graceful Fallbacks: Integrates fallback image URLs (Unsplash CDN) if TMDB poster paths are missing.\n• Network Caching: Employs CachedNetworkImage to eliminate duplicate HTTP requests and conserve bandwidth."
    papi2.font.size = Pt(14.5)
    papi2.font.color.rgb = TEXT_WHITE

    # ==========================================
    # SLIDE 10: Application Workflow (Priority Slide)
    # ==========================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide10)
    add_header(slide10, "Application User Journey & Workflow")

    # Workflow Diagram
    wf_steps = [
        {"title": "Splash Screen", "desc": "Animated startup & pre-loader"},
        {"title": "Main Navigation", "desc": "Material 3 bottom NavigationBar"},
        {"title": "Discover Feed", "desc": "Browse Home, Search, or Reels"},
        {"title": "Movie Details", "desc": "Read overview & view cast"},
        {"title": "Watchlist", "desc": "Bookmark favorite locally"}
    ]
    add_flow_diagram(slide10, Inches(0.8), Inches(1.65), Inches(11.733), Inches(2.35), wf_steps)

    # Workflow Highlights
    add_card(slide10, Inches(0.8), Inches(4.3), Inches(11.733), Inches(2.5))
    tb_wf = slide10.shapes.add_textbox(Inches(1.0), Inches(4.45), Inches(11.333), Inches(2.2))
    tf_wf = tb_wf.text_frame
    tf_wf.word_wrap = True

    pw1 = tf_wf.paragraphs[0]
    pw1.text = "Seamless User Experience Flow"
    pw1.font.size = Pt(20)
    pw1.font.bold = True
    pw1.font.color.rgb = GOLD

    pw2 = tf_wf.add_paragraph()
    pw2.text = "• Frictionless Browsing: Users land directly on the Home feed without complex login barriers.\n• State Preservation: IndexedStack preserves scroll position and search states across navigation tabs.\n• Instant Watchlist Updates: Adding a movie to Watchlist instantly updates the NavigationBar badge count."
    pw2.font.size = Pt(14.5)
    pw2.font.color.rgb = TEXT_WHITE

    # ==========================================
    # SLIDE 11: Future Scope (Priority Slide)
    # ==========================================
    slide11 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide11)
    add_header(slide11, "Future Scope & Roadmaps")

    futures = [
        {"title": "🔐 Cloud Authentication", "desc": "Firebase Auth / Supabase integration for multi-device sync and profiles."},
        {"title": "🤖 AI Recommendations", "desc": "Personalized movie suggestions tailored to user watchlist affinity."},
        {"title": "⏱️ Watch History", "desc": "Track trailer viewing history and resume video clips seamlessly."},
        {"title": "🔔 Push Notifications", "desc": "Real-time alerts when new trailers drop or bookmarked movies enter theaters."},
        {"title": "📥 Offline Caching", "desc": "Hive local database caching for viewing movie details without internet."},
        {"title": "📺 Video Streaming", "desc": "HLS / DASH video protocol integration for full-length legal movie streaming."}
    ]

    for idx, fut in enumerate(futures):
        col = idx % 2
        row = idx // 2
        left_pos = Inches(0.8 + col * 5.95)
        top_pos = Inches(1.65 + row * 1.75)

        add_card(slide11, left_pos, top_pos, Inches(5.75), Inches(1.6))
        tb_fs = slide11.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.15), Inches(5.35), Inches(1.3))
        tf_fs = tb_fs.text_frame
        tf_fs.word_wrap = True

        pfs1 = tf_fs.paragraphs[0]
        pfs1.text = fut["title"]
        pfs1.font.size = Pt(18)
        pfs1.font.bold = True
        pfs1.font.color.rgb = GOLD

        pfs2 = tf_fs.add_paragraph()
        pfs2.text = fut["desc"]
        pfs2.font.size = Pt(14)
        pfs2.font.color.rgb = TEXT_GRAY

    # ==========================================
    # SLIDE 12: Conclusion
    # ==========================================
    slide12 = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide12)
    add_header(slide12, "Conclusion & Project Summary")

    # Main Conclusion Card
    add_card(slide12, Inches(0.8), Inches(1.65), Inches(11.733), Inches(3.25), CARD_BG, CRIMSON)
    tb_c = slide12.shapes.add_textbox(Inches(1.1), Inches(1.85), Inches(11.133), Inches(2.85))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True

    pc1 = tf_c.paragraphs[0]
    pc1.text = "Cinemacentral Project Achievement Summary"
    pc1.font.size = Pt(24)
    pc1.font.bold = True
    pc1.font.color.rgb = TEXT_WHITE

    pc2 = tf_c.add_paragraph()
    pc2.text = "\nCinemacentral demonstrates modern mobile software engineering by delivering a high-performance, visually captivating movie discovery app built with Flutter and REST APIs."
    pc2.font.size = Pt(15)
    pc2.font.color.rgb = TEXT_WHITE

    pc3 = tf_c.add_paragraph()
    pc3.text = "\n• Flutter Mastery: Clean reactive UI, smooth Material 3 navigation, custom themes, and video components.\n• REST API Integration: Asynchronous HTTP fetching, JSON parsing, error handling, and lazy image loading.\n• UI/UX Excellence: Modern Netflix-inspired dark aesthetic suitable for production deployment."
    pc3.font.size = Pt(14)
    pc3.font.color.rgb = TEXT_GRAY

    # Thank You Banner
    add_card(slide12, Inches(0.8), Inches(5.1), Inches(11.733), Inches(1.7), CARD_BG, GOLD)
    tb_ty = slide12.shapes.add_textbox(Inches(1.0), Inches(5.25), Inches(11.333), Inches(1.4))
    tf_ty = tb_ty.text_frame
    tf_ty.word_wrap = True
    
    pty1 = tf_ty.paragraphs[0]
    pty1.text = "THANK YOU!"
    pty1.font.size = Pt(26)
    pty1.font.bold = True
    pty1.font.color.rgb = GOLD
    pty1.alignment = PP_ALIGN.CENTER

    pty2 = tf_ty.add_paragraph()
    pty2.text = "Questions & Answers  |  Cinemacentral Final Project Viva Presentation"
    pty2.font.size = Pt(16)
    pty2.font.color.rgb = TEXT_WHITE
    pty2.alignment = PP_ALIGN.CENTER

    # Save Presentation with fallback if file is locked by PowerPoint
    output_path = "Cinemacentral_Presentation.pptx"
    output_path_alt = "Cinemacentral_Presentation_Optimized.pptx"
    try:
        prs.save(output_path)
        print(f"SUCCESS: Created presentation at {os.path.abspath(output_path)}")
    except PermissionError:
        prs.save(output_path_alt)
        print(f"SUCCESS: Saved optimized presentation to {os.path.abspath(output_path_alt)} (as primary PPTX was locked by PowerPoint)")

if __name__ == "__main__":
    create_presentation()
